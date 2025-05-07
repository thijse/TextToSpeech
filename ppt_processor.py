"""
PowerPoint Processor Module

This module provides functionality to process PowerPoint presentations and convert them
to a Markdown format suitable for text-to-speech conversion.
"""

import os
import re
from pptx import Presentation
from typing import Optional, List, Dict


class PowerPointProcessor:
    """
    A processor for PowerPoint presentations that extracts notes and converts them
    to a Markdown format suitable for text-to-speech conversion.
    """
    
    def __init__(self):
        """
        Initialize the PowerPoint processor.
        """
        pass
    
    def _get_slide_title(self, slide) -> str:
        """
        Extract the title from a slide if available.
        
        Args:
            slide: A slide object from the python-pptx library.
            
        Returns:
            str: The slide title or empty string if no title is found.
        """
        try:
            # Look for a title placeholder
            for shape in slide.shapes:
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    if hasattr(shape, 'placeholder_format') and hasattr(shape.placeholder_format, 'type'):
                        if shape.placeholder_format.type == 1:  # Title placeholder
                            if hasattr(shape, 'text_frame') and shape.text_frame.text:
                                return self._sanitize_text(shape.text_frame.text)
            
            # If no title placeholder found, look for any text in the top of the slide
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame.text:
                    # Just take the first text we find as a fallback
                    return self._sanitize_text(shape.text_frame.text)
        except Exception as e:
            print(f"Warning: Could not extract title from slide: {e}")
        
        return ""
    
    def process_presentation(self, ppt_path: str, default_voice_name: str = None, 
                           include_empty_notes: bool = False, include_slide_titles: bool = True) -> str:
        """
        Process a PowerPoint presentation and convert it to a Markdown document.
        
        Args:
            ppt_path (str): The path to the PowerPoint file.
            default_voice_name (str, optional): The default voice name to use for the first section.
                                              Defaults to None.
            include_empty_notes (bool, optional): Whether to include slides with empty notes.
                                                Defaults to False.
            include_slide_titles (bool, optional): Whether to include slide titles in section headers.
                                                 Defaults to True.
            
        Returns:
            str: A Markdown document with the presentation content.
        """
        try:
            # Load the presentation
            presentation = Presentation(ppt_path)
            
            # Get the filename without extension for the main header
            filename = os.path.basename(ppt_path)
            filename_without_ext = os.path.splitext(filename)[0]
            
            # Start building the Markdown document
            markdown_lines = [f"# {filename_without_ext}\n"]
            
            # Process each slide
            first_slide_with_notes = True
            for i, slide in enumerate(presentation.slides):
                slide_number = i + 1
                
                # Extract notes text
                notes_text = self._extract_notes_text(slide)
                
                # Skip slides with empty notes if specified
                if not notes_text and not include_empty_notes:
                    continue
                
                # Get slide title if available
                slide_title = self._get_slide_title(slide)
                
                # Add slide as a subsection with title if available and enabled
                if slide_title and include_slide_titles:
                    section_header = f"## Slide {slide_number} - {slide_title}"
                else:
                    section_header = f"## Slide {slide_number}"
                
                markdown_lines.append(section_header + "\n")
                
                # Add voice annotation to the first slide with notes
                if first_slide_with_notes and default_voice_name:
                    markdown_lines.append(f"[voice:{default_voice_name}]\n")
                    first_slide_with_notes = False
                
                # Add notes text or a placeholder
                if notes_text:
                    markdown_lines.append(notes_text + "\n")
                #else:
                # No text in here. We do not want to have a lot non useful speech files   
                # markdown_lines.append("(No notes for this slide)\n")
            
            # Combine all lines into a single string
            return "\n".join(markdown_lines)
        
        except Exception as e:
            print(f"Error processing PowerPoint file: {e}")
            return ""
    
    def _sanitize_text(self, text: str) -> str:
        """
        Sanitize text by removing or replacing special characters.
        
        Args:
            text (str): The text to sanitize.
            
        Returns:
            str: The sanitized text.
        """
        if not text:
            return ""
            
        # Replace common problematic characters
        sanitized = text
        
        # Remove control characters (except newlines and tabs)
        sanitized = ''.join(ch for ch in sanitized if ch == '\n' or ch == '\t' or ord(ch) >= 32)
        
        # Replace non-breaking spaces with regular spaces
        sanitized = sanitized.replace('\xa0', ' ')
        
        # Remove zero-width spaces and other invisible characters
        sanitized = sanitized.replace('\u200b', '')  # zero-width space
        sanitized = sanitized.replace('\u200c', '')  # zero-width non-joiner
        sanitized = sanitized.replace('\u200d', '')  # zero-width joiner
        sanitized = sanitized.replace('\u2060', '')  # word joiner
        sanitized = sanitized.replace('\ufeff', '')  # zero-width no-break space
        
        # Replace other common special characters
        sanitized = sanitized.replace('\u2642', '')  # male sign (♂)
        sanitized = sanitized.replace('\u2640', '')  # female sign (♀)
        
        # Normalize whitespace (replace multiple spaces with a single space)
        sanitized = ' '.join(sanitized.split())
        
        return sanitized
    
    def _extract_notes_text(self, slide) -> str:
        """
        Extract notes text from a slide and sanitize it.
        
        Args:
            slide: A slide object from the python-pptx library.
            
        Returns:
            str: The sanitized notes text from the slide.
        """
        notes_slide = slide.notes_slide
        
        if notes_slide and notes_slide.notes_text_frame:
            raw_text = notes_slide.notes_text_frame.text.strip()
            return self._sanitize_text(raw_text)
        
        return ""
    
    def save_markdown(self, markdown_text: str, output_path: str) -> bool:
        """
        Save the Markdown document to a file.
        
        Args:
            markdown_text (str): The Markdown document text.
            output_path (str): The path where the Markdown file will be saved.
            
        Returns:
            bool: True if successful, False otherwise.
        """
        try:
            # Ensure the directory exists
            os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
            
            # Write the Markdown document to the file
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(markdown_text)
            
            return True
        except Exception as e:
            print(f"Error saving Markdown file: {e}")
            return False
